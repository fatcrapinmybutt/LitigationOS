"""
Text Extractor — Unified text extraction from any document type.
================================================================

Handles: PDF, DOCX, DOC, TXT, MD, HTML, CSV, XLSX, images (OCR).
Returns structured ExtractionResult with text, metadata, and page info.

Case-agnostic: extracts text without knowing what case it belongs to.
"""

import hashlib
import logging
import os
import re
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)


@dataclass
class PageText:
    """Text content from a single page or section."""
    page_number: int
    text: str
    char_count: int = 0

    def __post_init__(self):
        self.char_count = len(self.text)


@dataclass
class ExtractionResult:
    """Result of text extraction from a single file."""
    file_path: str
    file_name: str
    file_ext: str
    file_size: int
    sha256: str
    full_text: str
    pages: list[PageText] = field(default_factory=list)
    page_count: int = 0
    char_count: int = 0
    extraction_method: str = ""
    extracted_at: str = ""
    metadata: dict = field(default_factory=dict)
    error: Optional[str] = None

    def __post_init__(self):
        self.char_count = len(self.full_text)
        self.page_count = len(self.pages)


class TextExtractor:
    """Extracts text from any supported document format.

    Priority order:
      1. Native text extraction (PDF, DOCX, etc.)
      2. OCR fallback for images and image-PDFs
      3. Raw read for plain text formats
    """

    NATIVE_TEXT = {".txt", ".md", ".csv", ".html", ".htm", ".xml", ".json", ".rtf"}
    PDF_TYPES = {".pdf"}
    DOCX_TYPES = {".docx"}
    IMAGE_TYPES = {".jpg", ".jpeg", ".png", ".tiff", ".tif", ".bmp"}
    SPREADSHEET_TYPES = {".xlsx", ".xls"}
    DOC_TYPES = {".doc"}
    MSG_TYPES = {".msg"}
    EMAIL_TYPES = {".eml"}
    ODT_TYPES = {".odt"}
    PPTX_TYPES = {".pptx", ".ppt"}

    def __init__(self, ocr_enabled: bool = True):
        self.ocr_enabled = ocr_enabled
        # Lazy dependency flags — resolved on first use, not at construction
        self._dep_checked = False
        self._has_pymupdf = False
        self._has_docx = False
        self._has_tesseract = False
        self._has_openpyxl = False
        self._has_extract_msg = False
        self._has_odfpy = False
        self._has_pptx = False

    def _ensure_dependencies(self):
        """Lazy dependency check — called only when actually processing files.

        Avoids segfaults from fitz/PyMuPDF loading in subprocess contexts
        (exec_python, MCP command-runner) where native C libs may fail.
        """
        if self._dep_checked:
            return
        self._dep_checked = True

        # Each import is isolated — a segfault in one doesn't block others
        for lib, attr in [
            ("fitz", "_has_pymupdf"),
            ("docx", "_has_docx"),
            ("pytesseract", "_has_tesseract"),
            ("openpyxl", "_has_openpyxl"),
            ("extract_msg", "_has_extract_msg"),
            ("odf", "_has_odfpy"),
            ("pptx", "_has_pptx"),
        ]:
            try:
                __import__(lib)
                setattr(self, attr, True)
            except (ImportError, OSError):
                pass

    @property
    def SUPPORTED_EXTENSIONS(self) -> set[str]:
        """All file extensions this extractor can handle."""
        return (self.NATIVE_TEXT | self.PDF_TYPES | self.DOCX_TYPES |
                self.IMAGE_TYPES | self.SPREADSHEET_TYPES |
                self.DOC_TYPES | self.MSG_TYPES | self.EMAIL_TYPES |
                self.ODT_TYPES | self.PPTX_TYPES)

    def extract(self, file_path: str | Path) -> ExtractionResult:
        """Extract text from a file. Main entry point."""
        self._ensure_dependencies()
        path = Path(file_path).resolve()
        if not path.exists():
            return ExtractionResult(
                file_path=str(path),
                file_name=path.name,
                file_ext=path.suffix.lower(),
                file_size=0,
                sha256="",
                full_text="",
                error=f"File not found: {path}",
            )

        ext = path.suffix.lower()
        size = path.stat().st_size
        sha = self._hash_file(path)

        result = ExtractionResult(
            file_path=str(path),
            file_name=path.name,
            file_ext=ext,
            file_size=size,
            sha256=sha,
            full_text="",
            extracted_at=datetime.utcnow().isoformat(),
        )

        try:
            if ext in self.NATIVE_TEXT:
                self._extract_text_file(path, result)
            elif ext in self.PDF_TYPES:
                self._extract_pdf(path, result)
            elif ext in self.DOCX_TYPES:
                self._extract_docx(path, result)
            elif ext in self.IMAGE_TYPES:
                self._extract_image(path, result)
            elif ext in self.SPREADSHEET_TYPES:
                self._extract_spreadsheet(path, result)
            elif ext in self.DOC_TYPES:
                self._extract_doc(path, result)
            elif ext in self.MSG_TYPES:
                self._extract_msg(path, result)
            elif ext in self.EMAIL_TYPES:
                self._extract_eml(path, result)
            elif ext in self.ODT_TYPES:
                self._extract_odt(path, result)
            elif ext in self.PPTX_TYPES:
                self._extract_pptx(path, result)
            else:
                result.error = f"Unsupported file type: {ext}"
                result.extraction_method = "none"
        except Exception as e:
            result.error = f"Extraction failed: {type(e).__name__}: {e}"

        # Update computed fields after extraction methods populate full_text/pages
        result.char_count = len(result.full_text)
        result.page_count = len(result.pages)

        return result

    def _extract_text_file(self, path: Path, result: ExtractionResult):
        """Extract from plain text files."""
        for encoding in ("utf-8", "utf-8-sig", "cp1252", "latin-1"):
            try:
                text = path.read_text(encoding=encoding)
                result.full_text = text
                result.pages = [PageText(page_number=1, text=text)]
                result.extraction_method = f"text_read({encoding})"
                return
            except (UnicodeDecodeError, UnicodeError):
                continue
        result.error = "Could not decode text file with any encoding"

    def _extract_pdf(self, path: Path, result: ExtractionResult):
        """Extract from PDF using PyMuPDF with OCR fallback."""
        if not self._has_pymupdf:
            result.error = "PyMuPDF not installed: pip install pymupdf"
            return

        import fitz

        doc = fitz.open(str(path))
        result.metadata = dict(doc.metadata) if doc.metadata else {}
        pages = []
        all_text = []

        for i, page in enumerate(doc):
            try:
                text = page.get_text("text")
                # If page has very little text, might be image-based
                if len(text.strip()) < 20 and self.ocr_enabled and self._has_tesseract:
                    text = self._ocr_pdf_page(page) or text
                pages.append(PageText(page_number=i + 1, text=text))
                all_text.append(text)
            except Exception:
                pages.append(PageText(page_number=i + 1, text="[extraction error]"))

        doc.close()
        result.pages = pages
        result.full_text = "\n\n".join(all_text)
        result.extraction_method = "pymupdf" + ("+ocr" if self._has_tesseract else "")

    def _ocr_pdf_page(self, page) -> Optional[str]:
        """OCR a single PDF page via Tesseract."""
        try:
            import pytesseract
            from PIL import Image
            import io

            pix = page.get_pixmap(dpi=300)
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img)
            return text.strip() if text.strip() else None
        except Exception:
            return None

    def _extract_docx(self, path: Path, result: ExtractionResult):
        """Extract from DOCX using python-docx."""
        if not self._has_docx:
            result.error = "python-docx not installed: pip install python-docx"
            return

        import docx

        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        result.full_text = "\n".join(paragraphs)
        result.pages = [PageText(page_number=1, text=result.full_text)]
        result.extraction_method = "python-docx"

    def _extract_image(self, path: Path, result: ExtractionResult):
        """Extract text from image via OCR."""
        if not self._has_tesseract:
            result.error = "Tesseract not installed for OCR"
            return

        import pytesseract
        from PIL import Image

        img = Image.open(str(path))
        text = pytesseract.image_to_string(img)
        result.full_text = text
        result.pages = [PageText(page_number=1, text=text)]
        result.extraction_method = "tesseract_ocr"

    def _extract_spreadsheet(self, path: Path, result: ExtractionResult):
        """Extract from Excel spreadsheet."""
        if not self._has_openpyxl:
            result.error = "openpyxl not installed: pip install openpyxl"
            return

        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        all_text = []
        pages = []
        for i, sheet_name in enumerate(wb.sheetnames):
            sheet = wb[sheet_name]
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                rows.append(" | ".join(cells))
            sheet_text = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
            pages.append(PageText(page_number=i + 1, text=sheet_text))
            all_text.append(sheet_text)
        wb.close()

        result.full_text = "\n\n".join(all_text)
        result.pages = pages
        result.extraction_method = "openpyxl"

    def _extract_doc(self, path: Path, result: ExtractionResult):
        """Extract from legacy .doc using docx2txt or binary string fallback."""
        # Try docx2txt first (handles many .doc files despite the name)
        try:
            import docx2txt
            text = docx2txt.process(str(path))
            if text and text.strip():
                result.full_text = text
                result.pages = [PageText(page_number=1, text=text)]
                result.extraction_method = "docx2txt"
                return
        except Exception:
            pass
        # Fallback: extract printable ASCII runs from raw binary
        try:
            raw = path.read_bytes()
            text = raw.decode("utf-8", errors="ignore")
            import re as _re
            runs = _re.findall(r'[\x20-\x7E]{4,}', text)
            text = "\n".join(runs)
            result.full_text = text
            result.pages = [PageText(page_number=1, text=text)]
            result.extraction_method = "binary_strings"
        except Exception as e:
            result.error = f"DOC extraction failed: {e}"

    def _extract_msg(self, path: Path, result: ExtractionResult):
        """Extract from Outlook .msg files using extract-msg."""
        if not self._has_extract_msg:
            result.error = "extract-msg not installed: pip install extract-msg"
            return
        import extract_msg
        msg = extract_msg.Message(str(path))
        parts = []
        if msg.subject:
            parts.append(f"Subject: {msg.subject}")
        if msg.sender:
            parts.append(f"From: {msg.sender}")
        if msg.date:
            parts.append(f"Date: {msg.date}")
        if msg.body:
            parts.append(f"\n{msg.body}")
        text = "\n".join(parts)
        result.full_text = text
        result.pages = [PageText(page_number=1, text=text)]
        result.extraction_method = "extract_msg"
        result.metadata["subject"] = msg.subject or ""
        result.metadata["sender"] = msg.sender or ""
        result.metadata["date"] = str(msg.date) if msg.date else ""
        result.metadata["attachments"] = [a.longFilename for a in msg.attachments if a.longFilename]

    def _extract_eml(self, path: Path, result: ExtractionResult):
        """Extract from .eml files using stdlib email parser."""
        from email import policy
        from email.parser import BytesParser
        with open(path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)
        parts = []
        parts.append(f"Subject: {msg.get('subject', '')}")
        parts.append(f"From: {msg.get('from', '')}")
        parts.append(f"To: {msg.get('to', '')}")
        parts.append(f"Date: {msg.get('date', '')}")
        body = msg.get_body(preferencelist=("plain", "html"))
        if body:
            content = body.get_content()
            if content:
                parts.append(f"\n{content}")
        text = "\n".join(parts)
        result.full_text = text
        result.pages = [PageText(page_number=1, text=text)]
        result.extraction_method = "stdlib_email"
        result.metadata["subject"] = msg.get("subject", "")
        result.metadata["from"] = msg.get("from", "")
        result.metadata["to"] = msg.get("to", "")

    def _extract_odt(self, path: Path, result: ExtractionResult):
        """Extract from OpenDocument .odt files using odfpy."""
        if not self._has_odfpy:
            result.error = "odfpy not installed: pip install odfpy"
            return
        from odf.opendocument import load
        from odf.text import P
        doc = load(str(path))
        paragraphs = doc.getElementsByType(P)
        texts = []
        for p in paragraphs:
            t = ""
            for node in p.childNodes:
                if hasattr(node, "data"):
                    t += node.data
                elif hasattr(node, "__str__"):
                    t += str(node)
            if t:
                texts.append(t)
        text = "\n".join(texts)
        result.full_text = text
        result.pages = [PageText(page_number=1, text=text)]
        result.extraction_method = "odfpy"

    def _extract_pptx(self, path: Path, result: ExtractionResult):
        """Extract text from PowerPoint .pptx files using python-pptx."""
        if not self._has_pptx:
            result.error = "python-pptx not installed: pip install python-pptx"
            return
        from pptx import Presentation
        prs = Presentation(str(path))
        pages = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texts.append(para.text)
            slide_text = "\n".join(texts)
            pages.append(PageText(page_number=slide_num, text=slide_text))
        result.full_text = "\n\n".join(p.text for p in pages)
        result.pages = pages
        result.extraction_method = "python_pptx"

    @staticmethod
    def _hash_file(path: Path, block_size: int = 1048576) -> str:
        """Compute SHA-256 hash of file."""
        h = hashlib.sha256()
        with open(path, "rb") as f:
            while True:
                data = f.read(block_size)
                if not data:
                    break
                h.update(data)
        return h.hexdigest()
